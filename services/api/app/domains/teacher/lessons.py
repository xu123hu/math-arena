"""M3 教师端：教案 / 课件 / 讲解（§5.3 F1/F2/F11）。

- 教案从"复用并适配"开始：本地确定性模板 + 已有材料 + 规则建议（engine=local）；
  星辰可用时经 capability_gateway 走 wf_lesson_plan，失败仍可本地降级；
- 成品均是 draft Artifact，确认前不产生正式业务影响；
- 应用洞察：只改/新增 draft 版本，不原地篡改已确认版本。
"""

from __future__ import annotations

import io
import uuid

from fastapi import status
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.domains.teacher.artifacts import (
    _serialize_artifact,
    create_artifact,
    get_owned_artifact,
    update_artifact,
)
from app.domains.teacher.scope import assert_teacher_in_class, raise_http
from app.models.teacher import TeachingArtifact

ERR_NOT_FOUND = 40400

def _local_lesson_payload(topic: str, *, requirements: str | None, duration_minutes: int | None) -> dict:
    # 与 capability gateway 的本地降级模板共用，避免 provider 返回为空时退回旧占位时间线。
    from app.domains.teacher.capability_gateway import _local_lesson

    return {**_local_lesson(topic, requirements, duration_minutes), "insert_answer": []}


async def adapt_lesson(
    db: AsyncSession,
    teacher_id: uuid.UUID,
    class_id: uuid.UUID,
    *,
    topic: str,
    source_artifact_id: uuid.UUID | None,
    source_refs: list[str],
    requirements: str | None,
    duration_minutes: int | None,
    client_request_id: str,
) -> dict:
    await assert_teacher_in_class(db, teacher_id, class_id)

    import app.domains.teacher.capability_gateway as gw

    # 星辰可接入时走 wf_lesson_plan；其失败/未配置均由本地模板降级
    result = await gw.run_capability(
        db,
        teacher_id,
        scene="teacher.prep",
        class_id=class_id,
        capability="adapt_lesson",
        payload={
            "topic": topic,
            "requirements": requirements,
            "duration_minutes": duration_minutes,
        },
    )

    payload = result.get("payload") or _local_lesson_payload(topic, requirements=requirements, duration_minutes=duration_minutes)

    artifact = await create_artifact(
        db,
        owner_id=teacher_id,
        artifact_type="lesson_plan",
        scene="teacher.prep",
        class_id=class_id,
        payload=payload,
        source_refs=source_refs,
        engine=result.get("engine", "local"),
        degraded=result.get("degraded", False),
        warnings=(result.get("warnings") or []) if not result.get("degraded") else result.get("warnings") or ["本地模板模式"],
        parent_artifact_id=source_artifact_id,
        validation=result.get("validation", {"kind": "local_template", "deterministic": True}),
    )
    db.add(artifact)
    await db.flush()
    # 对齐前端 TeacherArtifact：返回完整 Artifact（content 含教案草稿）
    return _serialize_artifact(artifact)


async def create_lesson(
    db: AsyncSession,
    teacher_id: uuid.UUID,
    class_id: uuid.UUID,
    *,
    title: str,
    topic: str | None,
    duration_minutes: int | None,
    content: dict,
) -> dict:
    await assert_teacher_in_class(db, teacher_id, class_id)
    artifact = await create_artifact(
        db,
        owner_id=teacher_id,
        artifact_type="lesson_plan",
        scene="teacher.prep",
        class_id=class_id,
        payload={**content, "title": title, "topic": topic or title, "duration_minutes": duration_minutes},
        engine="local",
        degraded=False,
        validation={"kind": "manual", "deterministic": True},
    )
    db.add(artifact)
    await db.flush()
    return _serialize_artifact(artifact)


async def list_lessons(
    db: AsyncSession, teacher_id: uuid.UUID, class_id: uuid.UUID | None
) -> list[dict]:
    stmt = select(TeachingArtifact).where(
        TeachingArtifact.owner_id == teacher_id,
        TeachingArtifact.artifact_type == "lesson_plan",
        TeachingArtifact.deleted_at.is_(None),
    )
    if class_id:
        stmt = stmt.where(TeachingArtifact.class_id == class_id)
    rows = (await db.execute(stmt.order_by(TeachingArtifact.updated_at.desc()))).scalars().all()
    return [_serialize_artifact(a) for a in rows]


async def get_lesson(
    db: AsyncSession, teacher_id: uuid.UUID, lesson_id: uuid.UUID
) -> dict:
    a = await get_owned_artifact(db, teacher_id, lesson_id)
    if a.artifact_type != "lesson_plan":
        raise_http(ERR_NOT_FOUND, status.HTTP_404_NOT_FOUND, "not_found", recoverable=False)
    return {
        "lesson_id": str(a.id),
        "status": a.status,
        "version": a.version,
        "content": a.payload,
        "degraded": a.degraded,
        "engine": a.engine,
        "warnings": a.warnings or [],
    }


async def create_slides(
    db: AsyncSession,
    teacher_id: uuid.UUID,
    lesson_id: uuid.UUID,
    *,
    version: int,
    style: str | None,
    requirements: str | None,
) -> dict:
    lesson = await get_owned_artifact(db, teacher_id, lesson_id)
    # slides 是可发布产物：基于已确认教案生成 slide_deck draft
    if lesson.artifact_type != "lesson_plan":
        raise_http(ERR_NOT_FOUND, status.HTTP_404_NOT_FOUND, "not_found", recoverable=False)
    if lesson.status != "confirmed":
        raise_http(42210, 422, "confirmation_required", recoverable=True)
    if lesson.version != version:
        raise_http(40901, status.HTTP_409_CONFLICT, "version_conflict",
                   current_version=lesson.version, recoverable=True)
    payload = {
        "slides": _outline_from_lesson(lesson),
        "style": style or "课堂模板",
        "requirements": requirements,
        "topic": (lesson.payload or {}).get("topic") or "课堂教学",
        "objectives": (lesson.payload or {}).get("objectives") or [],
    }
    artifact = await create_artifact(
        db,
        owner_id=teacher_id,
        artifact_type="slide_deck",
        scene="teacher.prep",
        class_id=lesson.class_id,
        payload=payload,
        source_refs=[f"artifact:{lesson_id}"],
        engine="local",
        degraded=True,
        parent_artifact_id=lesson_id,
        validation={"kind": "outline", "deterministic": True},
    )
    db.add(artifact)
    artifact.payload = {
        **payload,
        "download_url": f"/api/teacher/slides/{artifact.id}/download",
        "filename": f"{payload['topic']}.pptx",
    }
    await db.flush()
    # 对齐前端契约：返回完整 Artifact（content.slides 为可编辑大纲）
    return _serialize_artifact(artifact)


async def render_slides_pptx(
    db: AsyncSession, teacher_id: uuid.UUID, slide_id: uuid.UUID
) -> tuple[bytes, str]:
    """把持久化 slide_deck artifact 确定性渲染为真正的 PPTX。"""
    artifact = await get_owned_artifact(db, teacher_id, slide_id)
    if artifact.artifact_type != "slide_deck":
        raise_http(ERR_NOT_FOUND, status.HTTP_404_NOT_FOUND, "not_found", recoverable=False)
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    payload = artifact.payload or {}
    topic = str(payload.get("topic") or "课堂教学")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic
    title_slide.placeholders[1].text = f"{payload.get('style') or '课堂模板'} · 本地可用课件"

    objectives = [str(x) for x in payload.get("objectives") or []]
    for index, spec in enumerate(payload.get("slides") or [], start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(spec.get("title") or f"第 {index} 页")
        frame = slide.placeholders[1].text_frame
        frame.clear()
        bullets = [str(x) for x in spec.get("bullets") or []]
        if not bullets:
            bullets = [f"教学环节：{spec.get('title') or '课堂'}"]
            if spec.get("notes"):
                bullets.append(f"建议用时：{spec['notes']}")
            if index == 1 and objectives:
                bullets.extend(f"教学目标：{objective}" for objective in objectives)
        for bullet_index, bullet in enumerate(bullets):
            paragraph = frame.paragraphs[0] if bullet_index == 0 else frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(24)
            paragraph.font.color.rgb = RGBColor(31, 41, 55)

    output = io.BytesIO()
    prs.save(output)
    filename = str(payload.get("filename") or f"{topic}.pptx")
    return output.getvalue(), filename


async def render_lesson_docx(
    db: AsyncSession, teacher_id: uuid.UUID, lesson_id: uuid.UUID
) -> tuple[bytes, str]:
    """把持久化教案产物确定性渲染为真正的 DOCX 文件。"""
    artifact = await get_owned_artifact(db, teacher_id, lesson_id)
    if artifact.artifact_type != "lesson_plan":
        raise_http(ERR_NOT_FOUND, status.HTTP_404_NOT_FOUND, "not_found", recoverable=False)

    from docx import Document
    from docx.shared import Pt

    payload = artifact.payload or {}
    topic = str(payload.get("topic") or payload.get("title") or "课堂教案")
    document = Document()
    document.add_heading(topic, level=0)
    document.add_paragraph(f"课时：{payload.get('duration_minutes') or 45} 分钟")

    document.add_heading("教学目标", level=1)
    objectives = payload.get("objectives") or [f"掌握{topic}的核心概念与基本方法"]
    for objective in objectives:
        document.add_paragraph(str(objective), style="List Bullet")

    document.add_heading("教学过程", level=1)
    for index, step in enumerate(payload.get("timeline") or [], start=1):
        phase = str(step.get("phase") or f"环节 {index}")
        minutes = step.get("minutes")
        document.add_heading(f"{index}. {phase}" + (f"（{minutes} 分钟）" if minutes else ""), level=2)
        activities = step.get("activities") or []
        if activities:
            for activity in activities:
                document.add_paragraph(str(activity), style="List Bullet")
        else:
            document.add_paragraph("按班级实际情况组织讲授、提问与练习。")

    requirements = payload.get("requirements")
    if requirements:
        document.add_heading("备课要求", level=1)
        document.add_paragraph(str(requirements))

    styles = document.styles
    styles["Normal"].font.name = "Microsoft YaHei"
    styles["Normal"].font.size = Pt(11)
    output = io.BytesIO()
    document.save(output)
    return output.getvalue(), f"{topic}.docx"


def _outline_from_lesson(lesson: TeachingArtifact) -> list[dict]:
    timeline = (lesson.payload or {}).get("timeline") or []
    return [
        {"page": i + 1, "title": p.get("phase", f"第{i+1}页"),
         "bullets": [str(activity) for activity in p.get("activities", []) if str(activity).strip()],
         "notes": f"约{p.get('minutes', '')}分钟"}
        for i, p in enumerate(timeline)
    ] or [{"page": 1, "title": "课堂", "bullets": [], "notes": ""}]


async def create_explainer(
    db: AsyncSession,
    teacher_id: uuid.UUID,
    class_id: uuid.UUID | None,
    *,
    question: str,
    reference_solution: str | None,
    target_minutes: int | None,
) -> dict:
    """讲题卡（explanation draft）：class_id 可空（讲题可不绑定班级）。"""
    if class_id:
        await assert_teacher_in_class(db, teacher_id, class_id)
    artifact = await create_artifact(
        db,
        owner_id=teacher_id,
        artifact_type="explanation",
        scene="teacher.prep",
        class_id=class_id,
        payload={
            "question": question,
            "reference_solution": reference_solution,
            "target_minutes": target_minutes,
            "steps": _explanatory_steps(question),
        },
        engine="local",
        degraded=False,
        validation={"kind": "deterministic", "verified": False},
    )
    db.add(artifact)
    await db.flush()
    return _serialize_artifact(artifact)


def _explanatory_steps(question: str) -> list[dict]:
    return [
        {"step": 1, "instruction": "重述问题与关键条件", "detail": question, "is_certain": True},
        {"step": 2, "instruction": "给出分步推导方向", "detail": "需结合标准答案与教师材料", "is_certain": False},
        {"step": 3, "instruction": "易错点提醒", "detail": "待教师依据本班错因补充", "is_certain": False},
    ]


async def apply_insight_to_lesson(
    db: AsyncSession,
    teacher_id: uuid.UUID,
    lesson_id: uuid.UUID,
    *,
    insight_id: uuid.UUID,
    version: int,
    instruction: str | None,
) -> dict:
    """将可行动洞察写入教案草稿（审计 I-08：按 insight_id 加载真实洞察，不信任前端摘要）。"""
    from app.domains.teacher.scope import assert_teacher_in_class
    from app.models.teacher import ActionableInsight

    lesson = await get_owned_artifact(db, teacher_id, lesson_id)
    insight = await db.get(ActionableInsight, insight_id)
    if insight is None:
        raise_http(ERR_NOT_FOUND, status.HTTP_404_NOT_FOUND, "insight_not_found", recoverable=False)
    # 洞察班级必须与教师范围一致
    if lesson.class_id and insight.class_id != lesson.class_id:
        await assert_teacher_in_class(db, teacher_id, insight.class_id)
    if insight.status != "active":
        raise_http(40001, 422, "insight_not_active", recoverable=True)

    new_content = dict(lesson.payload or {})
    new_content["insert_answer"] = list(new_content.get("insert_answer", [])) + [
        {
            "kind": "insight",
            "insight_id": str(insight.id),
            "summary": insight.summary,
            "evidence": insight.evidence,
            "instruction": instruction,
        }
    ]
    updated, _created = await update_artifact(
        db, teacher_id, lesson_id, version=version, payload=new_content
    )
    # 洞察标记为 applied（幂等：重复应用由版本冲突拦截）
    insight.status = "applied"
    await db.flush()
    return _serialize_artifact(updated)
